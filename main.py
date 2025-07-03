
#%%
import subprocess
import matplotlib.ticker as ticker
import io
import os
import matplotlib.pyplot as plt
import pandas as pd 
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import dash
from pptx.util import Inches
from dash import dcc, html, no_update
from dash.dependencies import Input, Output, State

app = dash.Dash(__name__)
df1 = pd.read_csv('online_sales.csv', delimiter = ',')

# Function to get age categories from data
def get_age_categories(df1):
    if df1.empty:
        return []
    # Get min and max ages, rounded down to nearest 5
    min_age = (min(df1['age']) // 5) * 5
    max_age = (max(df1['age']) // 5) * 5
    
    # Create age ranges in steps of 5
    age_bins = range(min_age, max_age + 6, 5)
    
    # Build the age category labels
    age_labels = []
    for i in age_bins[:-1]:
        category = str(i) + "-" + str(i + 4)
        age_labels.append(category)
    
    # Convert to the expected format
    result = []
    for category in age_labels:
        result.append({'label': category, 'value': category})
    
    return result

# Function to define ppt layout and specifications (where KPIs are placed)
def set_custom_fill_and_outline(shape, is_large_rectangle = False):
    if is_large_rectangle:
        fill_color = RGBColor(244,244,244)
    else:
        fill_color = RGBColor(222,222,223)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color 
    shape.line.color.rgb = fill_color 

#Function to calculate KPIs 
def calculate_kpis(df1):
    if df1.empty: 
        return 0,0,0
    
    total_new_users = df1['new_user'].sum()
    total_converted = df1['converted'].sum()
    total_pages_visited = df1['total_pages_visited'].sum()
    total_conversion = df1['converted'].sum()
    conversion_rate = round((total_conversion/total_pages_visited)*100,2) if total_pages_visited > 0 else 0

    return total_new_users, total_converted, conversion_rate

#Function to add KPI metrics to slides
def add_kpi(slide, left, top, value, label):
    text_box = slide.shapes.add_textbox(left, top, Inches(2), Inches(1))
    text_frame = text_box.text_frame
    text_frame.clear()
    p = text_frame.paragraphs[0]
    run = p.add_run()
    run.text = value
    run.font.size = Pt(16)
    run.font.bold = True
    run.font.color.rgb = RGBColor(0, 51, 102)
    p = text_frame.add_paragraph()
    run = p.add_run()
    run.text = label
    run.font.size = Pt(11)
    run.font.color.rgb = RGBColor(0, 51, 102)

#Function to add the Key statistics heading
def add_heading_text(slide, left, top, text):
    textbox = slide.shapes.add_textbox(left, top, Inches(3), Inches(0.5))
    text_frame = textbox.text_frame
    text_frame.text = text
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.CENTER
    paragraph.font.size = Pt(18)
    paragraph.font.bold = True
    textbox.left = int(left - (textbox.width / 2))

#Function to clone shapes from source slide 
def clone_shapes(source_slide, new_slide):
    for shape in source_slide.shapes:
        if not shape.is_placeholder:
            if shape.shape_type == MSO_SHAPE.RECTANGLE:
                new_shape = new_slide.shapes.add_shape(
                    shape.auto_shape_type,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                is_large_rectangle = shape.width > Inches(5.5)
                set_custom_fill_and_outline(new_shape, is_large_rectangle)
                if hasattr(new_shape, 'text_frame'):
                    new_shape._element.remove(new_shape._element.txBody)

            elif shape.shape_type == MSO_SHAPE.ROUNDED_RECTANGLE:
                new_shape = new_slide.shapes.add_shape(
                    MSO_SHAPE.ROUNDED_RECTANGLE,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )
                is_large_rectangle = shape.width > Inches(5.5)
                set_custom_fill_and_outline(new_shape, is_large_rectangle)
                if hasattr(new_shape, 'text_frame'):
                    new_shape._element.remove(new_shape._element.txBody)

            elif shape.shape_type == 13: 
                image_stream = shape.image.blob
                image_file = io.BytesIO(image_stream)
                new_slide.shapes.add_picture(
                    image_file,
                    shape.left,
                    shape.top,
                    shape.width,
                    shape.height
                )

# Improved function to check if Inkscape is available
def is_inkscape_available():
    try:
        with open(os.devnull, 'w') as devnull:
            subprocess.run(["inkscape", "--version"], 
                         check=True, stdout=devnull, stderr=devnull)
        return True
    except (subprocess.CalledProcessError, FileNotFoundError):
        return False

# Improved function to convert SVG to EMF with fallback
def convert_svg_to_emf(svg_path, emf_path):
    if not os.path.exists(svg_path):
        print("SVG file not found:",svg_path)
        return None
        
    # Check if Inkscape is available for error handling
    if is_inkscape_available():
        try:
            with open(os.devnull, 'w') as devnull:
                subprocess.run(
                    ["inkscape", svg_path, "export-type=emf", "export-filename", emf_path],
                    check=True,
                    stdout=devnull,
                    stderr=devnull
                )
            if os.path.exists(emf_path):
                print("Successfully converted to EMF:", emf_path)
                return emf_path
        except subprocess.CalledProcessError as e:
            print("Error converting SVG to EMF:", e)
    
    # Fallback: convert to PNG using matplotlib
    print("Inkscape not available or conversion failed. Using PNG fallback")
    png_path = svg_path.replace('.svg', '.png')
    try:
        # Re-generate the plot as PNG directly
        return png_path 
    except Exception as e:
        print("Error in PNG fallback:",e)
        return None

def generate_total_sites_chart(df1, title_suffix=""):
    try:
        # Check if we have any data to work with
        if df1.empty:
            print("DataFrame is empty, cannot generate chart")
            return None
        
        # Calculate age ranges 
        min_age = (min(df1['age']) // 5) * 5
        max_age = (max(df1['age']) // 5) * 5
        age_bins = range(min_age, max_age + 6, 5)
        age_labels = []
        for i in age_bins[:-1]:
            label = str(i) + "-" + str(i + 4)
            age_labels.append(label)
        
        # Work with a copy to avoid modifying the original dataframe
        df1_copy = df1.copy()
        
        # Create age groups using pandas cut function
        df1_copy['AgeGroup'] = pd.cut(df1_copy['age'], bins=age_bins, labels=age_labels, right=False)
        
        # Group by age and sum up total pages visited
        age_group_stats = df1_copy.groupby('AgeGroup')['total_pages_visited'].sum().reset_index()
        
        # Set up the plot
        plt.figure(figsize=(10, 6))
        bars = plt.bar(age_group_stats['AgeGroup'], age_group_stats['total_pages_visited'],
                       color='skyblue', edgecolor='black')
        
        # Add numbers on top of each bar
        for bar in bars:
            height = bar.get_height()
            x_pos = bar.get_x() + bar.get_width() / 2.0
            plt.text(x_pos, height, str(int(height)), ha='center', va='bottom')
        
        # Set up the title
        if title_suffix:
            title = "Total pages visited vs Age group " + title_suffix
        else:
            title = "Total pages visited vs Age group"
        title = title.strip()
        
        # Customize the chart appearance
        plt.title(title, fontsize=14, pad=20)
        plt.xlabel('Age Group (5-year ranges)', fontsize=12)
        plt.ylabel('Total Pages Visited', fontsize=12)
        plt.grid(axis='y', linestyle='', alpha=0.7)
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        # Remove chart borders/spines
        ax = plt.gca()
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['bottom'].set_visible(True)  # Keep x-axis visible
        
        # Set up file paths
        svg_path = 'temp_chart.svg'
        emf_path = 'output_chart.emf'
        png_path = 'output_chart.png'
        
        # Save as SVG first
        plt.savefig(svg_path, format='svg', bbox_inches='tight', dpi=600, transparent=True)
        
        # Try to convert SVG to EMF format
        result_path = convert_svg_to_emf(svg_path, emf_path)
        
        # Check if EMF conversion worked
        if result_path is None or not os.path.exists(result_path):
            # EMF didn't work, use PNG instead
            print("Using PNG fallback...")
            plt.savefig(png_path, format='png', bbox_inches='tight', dpi=300, transparent=True)
            plt.close()
            
            # Clean up the temporary SVG file
            if os.path.exists(svg_path):
                os.remove(svg_path)
            
            # Check if PNG was created successfully
            if os.path.exists(png_path):
                print("Chart saved as PNG: " + png_path)
                return png_path
            else:
                print("Failed to save PNG chart")
                return None
        
        # EMF worked, close the plot and clean up
        plt.close()
        print("Chart saved as EMF: ", result_path)
        return result_path
    
    except Exception as e:
        print("Error generating chart: " + str(e))
        plt.close()  # Make sure to close plot even if something went wrong
        return None

def generate_conversion_chart(df1, title_suffix=""):
    try:
        # Check if we have data to work with
        if df1.empty:
            print("DataFrame is empty, cannot generate chart")
            return None
        
        # Calculate age ranges - round down to nearest 5
        min_age = (min(df1['age']) // 5) * 5
        max_age = (max(df1['age']) // 5) * 5
        age_bins = range(min_age, max_age + 6, 5)
        age_labels = []
        for i in age_bins[:-1]:
            label = str(i) + "-" + str(i + 4)
            age_labels.append(label)
        
        # Work with a copy to avoid modifying original dataframe
        df1_copy = df1.copy()
        df1_copy['AgeGroup'] = pd.cut(df1_copy['age'], bins=age_bins, labels=age_labels, right=False)
        
        # Calculate average conversion rate by age group
        age_group_stats = df1_copy.groupby('AgeGroup')['converted'].mean().reset_index()
        
        # Create the plot
        plt.figure(figsize=(10, 6))
        bars = plt.bar(age_group_stats['AgeGroup'], age_group_stats['converted'], 
                      color='#003060', edgecolor='black')
        
        # Set up the title
        if title_suffix:
            title = "Average conversion rate vs Age group " + title_suffix
        else:
            title = "Average conversion rate vs Age group"
        title = title.strip()
        
        # Customize the chart
        plt.title(title, fontsize=14, pad=20)
        plt.xlabel('Age Group (5-year ranges)', fontsize=12)
        plt.ylabel('Conversion Rate (%)', fontsize=12)
        plt.grid(axis='y', linestyle='', alpha=0.3)
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        # Remove chart borders
        ax = plt.gca()
        ax.spines['top'].set_visible(False)
        ax.spines['right'].set_visible(False)
        ax.spines['left'].set_visible(False)
        ax.spines['bottom'].set_visible(True)
        
        # Set up file paths
        svg_path = 'temp_conversion_chart.svg'
        emf_path = 'conversion_chart.emf'
        png_path = 'conversion_chart.png'
        
        # Save as SVG first
        plt.savefig(svg_path, format='svg', bbox_inches='tight', dpi=600, transparent=True)
        
        # Try to convert to EMF
        result_path = convert_svg_to_emf(svg_path, emf_path)
        
        # Check if EMF conversion worked
        if result_path is None or not os.path.exists(result_path):
            # Use PNG as fallback
            print("Using PNG fallback for conversion chart...")
            plt.savefig(png_path, format='png', bbox_inches='tight', dpi=300, transparent=True)
            plt.close()
            
            # Clean up temporary SVG file
            if os.path.exists(svg_path):
                os.remove(svg_path)
            
            # Check if PNG was created successfully
            if os.path.exists(png_path):
                print("Conversion chart saved as PNG: " + png_path)
                return png_path
            else:
                print("Failed to save PNG conversion chart")
                return None
        
        # EMF worked, clean up and return
        plt.close()
        print("Conversion chart saved as EMF: " + result_path)
        return result_path
    
    except Exception as e:
        print("Error generating conversion chart: " + str(e))
        plt.close() 
        return None


def generate_total_sites_chart(df1, title_suffix=""):
    try:
        # Check if dataframe has data
        if df1.empty:
            print("DataFrame is empty, cannot generate total sites chart")
            return None
        
        # Calculate age ranges
        min_age = (min(df1['age']) // 5) * 5
        max_age = (max(df1['age']) // 5) * 5
        age_bins = range(min_age, max_age + 6, 5)
        age_labels = []
        for i in age_bins[:-1]:
            label = str(i) + "-" + str(i + 4)
            age_labels.append(label)
        
        # Work with copy of dataframe
        df1_copy = df1.copy()
        df1_copy['AgeGroup'] = pd.cut(df1_copy['age'], bins=age_bins, labels=age_labels, right=False)
        
        # Count total users by age group
        age_group_counts = df1_copy.groupby('AgeGroup').size().reset_index(name='total_users')
        
        # Create the plot
        plt.figure(figsize=(10, 6))
        bars = plt.bar(age_group_counts['AgeGroup'], age_group_counts['total_users'], 
                      color='#00008B', edgecolor='black')
        
        # Build the title
        if title_suffix:
            title = "Total Users by Age Group " + title_suffix
        else:
            title = "Total Users by Age Group"
        title = title.strip()
        
        # Customize the plot
        plt.title(title, fontsize=14, pad=20)
        plt.xlabel('Age Group (5-year ranges)', fontsize=12)
        plt.ylabel('Total Users', fontsize=12)
        plt.grid(axis='y', linestyle='', alpha=0.3)
        plt.xticks(rotation=45, ha='right')
        plt.tight_layout()
        
        # Set up file paths
        svg_path = 'temp_total_sites_chart.svg'
        emf_path = 'total_sites_chart.emf'
        png_path = 'total_sites_chart.png'
        
        # Save as SVG first
        plt.savefig(svg_path, format='svg', bbox_inches='tight', dpi=600, transparent=True)
        
        # Try EMF conversion
        result_path = convert_svg_to_emf(svg_path, emf_path)
        
        # Check if EMF conversion succeeded
        if result_path is None or not os.path.exists(result_path):
            # Fallback to PNG
            print("Using PNG fallback for total sites chart...")
            plt.savefig(png_path, format='png', bbox_inches='tight', dpi=300, transparent=True)
            plt.close()
            
            # Remove temporary SVG
            if os.path.exists(svg_path):
                os.remove(svg_path)
            
            # Verify PNG was created
            if os.path.exists(png_path):
                print("Total sites chart saved as PNG: " , png_path)
                return png_path
            else:
                print("Failed to save PNG total sites chart")
                return None
        
        # EMF conversion successful
        plt.close()
        print("Total sites chart saved as EMF: " , result_path)
        return result_path
    
    except Exception as e:
        print("Error generating total sites chart: ",e)
        plt.close()
        return None

def add_charts_to_presentation(prs, df_filtered, slide_index=1):
    try:
        # Check if the slide index is valid
        if len(prs.slides) <= slide_index:
            error_msg = "Slide index " + str(slide_index) + " out of range (presentation has " + str(len(prs.slides)) + " slides)"
            raise ValueError(error_msg)
        
        # Calculate the KPI values
        print("Calculating KPIs...")
        total_new_users, total_converted, conversion_rate = calculate_kpis(df_filtered)
        print("KPIs calculated - New Users:", total_new_users, "Converted:", total_converted, "Rate:", conversion_rate, "%")
        
        # Get the slide we want to work with
        slide = prs.slides[slide_index]
        print("Working on slide", slide_index)
        
        # Update KPI placeholders
        print("Updating KPI placeholder")
        kpis_updated = False
        
        for shape in slide.shapes:
            if shape.has_text_frame:
                text = shape.text_frame.text.strip()
                
                if text == 'A':
                    shape.text_frame.text = "Total new users: " + str(total_new_users)
                    kpis_updated = True
                    print("Updated placeholder A with:", total_new_users)
                elif text == 'B':
                    shape.text_frame.text = "Total converted:" + str(total_converted)  
                    kpis_updated = True
                    print("Updated placeholder B with:", total_converted)
                elif text == 'C':
                    shape.text_frame.text = "Conversion rate:" + str(conversion_rate) + "%"  
                    kpis_updated = True
                    print("Updated placeholder C with:", conversion_rate, "%")
        
        # Check if we found any KPI placeholders
        if not kpis_updated:
            print("No KPI placeholders (A, B, C) were found on the input ppt slide - check the input the slide")
        
        # Generate both charts
        print("Generating charts.")
        chart_path = generate_total_sites_chart(df_filtered, title_suffix="")
        chart_path1 = generate_conversion_chart(df_filtered, title_suffix="")

        # Add the first chart if it was created successfully
        if chart_path and os.path.exists(chart_path):
            print("Total sites chart generated at:", chart_path)
            try:
                # Set position for first chart (left side)
                left = Inches(0.5)  
                top = Inches(3.0)
                width = Inches(6)   
                height = Inches(4)
                
                # Add the picture to the slide
                pic = slide.shapes.add_picture(chart_path, left, top, width, height)
                print("Total sites chart added to slide successfully")
                
                # Look for chart title placeholder and update it
                title_updated = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip() == 'D':
                        shape.text_frame.text = ""
                        title_updated = True
                        break
                
                if not title_updated:
                    print("Chart title placeholder 'D' not found")
                
            except Exception as e:
                print("Error adding total sites chart to slide:", e)
                return None
        else:
            print("Error: Total sites chart not generated or file not found")
            return None

        # Add the second chart if it was created successfully
        if chart_path1 and os.path.exists(chart_path1):
            print("Conversion chart generated at:", chart_path1)
            try:
                # Set position for second chart (right side)
                left = Inches(7)    
                top = Inches(3.0)
                width = Inches(6)   
                height = Inches(4)
                
                print("Adding conversion chart to slide at position: " + str(left.inches) + "in, " + str(top.inches) + "in")
                
                # Add the picture to the slide
                pic = slide.shapes.add_picture(chart_path1, left, top, width, height)
                print("Conversion chart added to slide successfully")
                
                # Look for chart title placeholder and update it
                title_updated = False
                for shape in slide.shapes:
                    if shape.has_text_frame and shape.text_frame.text.strip() == 'E':
                        shape.text_frame.text = ""
                        title_updated = True
                        break
                
                if not title_updated:
                    print("Chart title placeholder 'E' not found")
                
            except Exception as e:
                print("Error adding conversion chart to slide:", e)
                return None
        else:
            print("Error: Conversion chart not generated")
            return None
        
        print("Charts added successfully")
        
        # Clean up all temporary files
        temp_files = [
            'temp_conversion_chart.svg', 'conversion_chart.emf', 'conversion_chart.png',
            'temp_total_sites_chart.svg', 'total_sites_chart.emf', 'total_sites_chart.png'
        ]
        
        for temp_file in temp_files:
            if os.path.exists(temp_file):
                try:
                    os.remove(temp_file)
                    print("Removed temporary file:", temp_file)
                except Exception as e:
                    print("Warning: Could not remove", temp_file, ":", e)
        
        return prs
        
    except Exception as e:
        print("Critical error in add_charts_to_presentation:", e)
        import traceback
        traceback.print_exc()
        return None


def filter_dataframe_by_age(df1, selected_age_category):
    # Check if we should return all data
    if selected_age_category == 'all' or not selected_age_category:
        return df1
    
    # Parse the age range string (e.g., "25-29" becomes min=25, max=29)
    age_range = selected_age_category.split('-')
    min_age = int(age_range[0])
    max_age = int(age_range[1])
    
    # Filter the dataframe based on age range
    filtered_df = df1[(df1['age'] >= min_age) & (df1['age'] <= max_age)]
    return filtered_df


def create_presentation(df_filtered, template_path="Sales_presentation1.pptx"):
    try:
        # Check if template file exists
        if not os.path.exists(template_path):
            error_msg = "Template file not found at " + template_path
            raise FileNotFoundError(error_msg)
        
        # Load the template presentation
        prs = Presentation(template_path)
        
        # Make sure we have enough slides to work with
        if len(prs.slides) < 2:
            raise ValueError("Template must have at least 2 slides")
        
        # Update the title slide with current information
        title_slide = prs.slides[0]
        try:
            # Update title if it exists
            if hasattr(title_slide.shapes, 'title') and title_slide.shapes.title:
                title_slide.shapes.title.text = "Sales Dashboard Report"
            
            # Update subtitle with current date if placeholder exists
            if len(title_slide.placeholders) > 1:
                current_time = datetime.now().strftime('%Y-%m-%d %H:%M:%S')
                title_slide.placeholders[1].text = "Generated on " + current_time
                
        except Exception as e:
            print("Warning: Could not update title slide: " + str(e))
        
        # Add KPIs and charts to the second slide
        prs = add_charts_to_presentation(prs, df_filtered, slide_index=1)
        
        # Check if chart addition was successful
        if prs is None:
            raise ValueError("Failed to add charts to presentation")
        
        # Generate filename with timestamp
        current_timestamp = datetime.now().strftime('%Y%m%d_%H%M%S')
        ppt_filename = "sales_report_" + current_timestamp + ".pptx"
        
        # Save the presentation
        prs.save(ppt_filename)
        print("Presentation saved as: " + ppt_filename)
        
        return ppt_filename
        
    except Exception as e:
        print("Error creating presentation: " + str(e))
        import traceback
        traceback.print_exc()
        return None


# Create the Dash layout with professional styling
app.layout = html.Div([
    # Header Section
    html.Div([
        html.H1("Sales Department Dashboard", 
                style={
                    'color': '#0051a6', 
                    'textAlign': 'center', 
                    'marginBottom': '0px',
                    'fontSize': '2.5rem',
                    'fontWeight': '700',
                    'letterSpacing': '1px'
                }),
        html.Hr(style={
            'border': 'none',
            'height': '3px',
            'background': 'linear-gradient(90deg, #0051a6, #007bff)',
            'margin': '20px auto',
            'width': '300px',
            'borderRadius': '2px'
        })
    ], style={
        'backgroundColor': 'white',
        'padding': '40px 20px 30px 20px',
        'boxShadow': '0 2px 10px rgba(0,0,0,0.1)',
        'marginBottom': '40px'
    }),
    
    # Main Content Container
    html.Div([
        # Control Panel Section
        html.Div([
            html.H3("Dashboard Controls", 
                   style={
                       'color': '#0051a6', 
                       'marginBottom': '25px',
                       'fontSize': '1.4rem',
                       'fontWeight': '600'
                   }),
            
            html.Div([
                # Age Category Filter Section
                html.Div([
                    html.Label("Select Age Category:", 
                              style={
                                  'fontWeight': '600', 
                                  'marginBottom': '12px', 
                                  'color': '#333',
                                  'fontSize': '1rem',
                                  'display': 'block'
                              }),
                    dcc.Dropdown(
                        id='age-category-dropdown',
                        options=[{'label': 'All Age Groups', 'value': 'all'}] + get_age_categories(df1),
                        value='all',
                        multi=True,
                        placeholder="Choose age categories...",
                        style={
                            'width': '350px', 
                            'marginBottom': '10px',
                            'fontSize': '0.95rem'
                        }
                    ),
                    html.Small("Select multiple categories to compare different age groups", 
                             style={
                                 'color': '#666', 
                                 'fontStyle': 'italic',
                                 'fontSize': '0.85rem'
                             })
                ], style={
                    'display': 'inline-block', 
                    'marginRight': '80px', 
                    'verticalAlign': 'top',
                    'minWidth': '350px'
                }),
                
                # Download Section
                html.Div([
                    html.Label("Generate Report:", 
                              style={
                                  'fontWeight': '600', 
                                  'marginBottom': '12px', 
                                  'color': '#333',
                                  'fontSize': '1rem',
                                  'display': 'block'
                              }),
                    html.Button("📊 Download PowerPoint Report", 
                               id="download-btn", 
                               n_clicks=0,
                               style={
                                   'backgroundColor': '#0051a6',
                                   'color': 'white',
                                   'border': 'none',
                                   'padding': '12px 24px',
                                   'fontSize': '1rem',
                                   'fontWeight': '600',
                                   'borderRadius': '8px',
                                   'cursor': 'pointer',
                                   'transition': 'all 0.3s ease',
                                   'boxShadow': '0 4px 8px rgba(0,81,166,0.3)',
                                   'minWidth': '250px'
                               }),
                    dcc.Download(id="download-ppt"),
                    html.Small("Click to generate and download a comprehensive report", 
                             style={
                                 'color': '#666', 
                                 'fontStyle': 'italic',
                                 'fontSize': '0.85rem',
                                 'display': 'block',
                                 'marginTop': '8px'
                             })
                ], style={
                    'display': 'inline-block', 
                    'verticalAlign': 'top'
                }),
                
            ], style={
                'display': 'flex',
                'flexWrap': 'wrap',
                'alignItems': 'flex-start',
                'gap': '40px'
            }),
            
        ], style={
            'backgroundColor': '#f8f9fa',
            'padding': '35px',
            'borderRadius': '15px',
            'marginBottom': '50px',
            'boxShadow': '0 4px 12px rgba(0,0,0,0.08)',
            'border': '1px solid #e9ecef'
        }),
        
        # KPI Cards Section
        html.Div([
            html.H3("Key Performance Indicators", 
                   style={
                       'color': '#0051a6', 
                       'textAlign': 'center', 
                       'marginBottom': '35px',
                       'fontSize': '1.6rem',
                       'fontWeight': '600'
                   }),
            
            html.Div([
                # New Users Card
                html.Div([
                    html.Div([
                        html.I(className="fas fa-user-plus", style={
                            'fontSize': '2rem', 
                            'color': '#28a745',
                            'marginBottom': '15px'
                        }),
                        html.H4("Total New Users", style={
                            'color': '#333', 
                            'marginBottom': '12px',
                            'fontSize': '1.1rem',
                            'fontWeight': '600'
                        }),
                        html.H2(id="kpi-new-users", style={
                            'color': '#28a745', 
                            'margin': '0',
                            'fontSize': '2.2rem',
                            'fontWeight': '700'
                        })
                    ], style={'textAlign': 'center'})
                ], style={
                    'backgroundColor': 'white',
                    'padding': '30px 25px',
                    'borderRadius': '12px',
                    'boxShadow': '0 6px 20px rgba(40,167,69,0.15)',
                    'border': '1px solid #e9ecef',
                    'minWidth': '220px',
                    'transition': 'transform 0.3s ease, box-shadow 0.3s ease'
                }),
                
                # Converted Users Card
                html.Div([
                    html.Div([
                        html.I(className="fas fa-check-circle", style={
                            'fontSize': '2rem', 
                            'color': '#dc3545',
                            'marginBottom': '15px'
                        }),
                        html.H4("Total Converted", style={
                            'color': '#333', 
                            'marginBottom': '12px',
                            'fontSize': '1.1rem',
                            'fontWeight': '600'
                        }),
                        html.H2(id="kpi-converted", style={
                            'color': '#dc3545', 
                            'margin': '0',
                            'fontSize': '2.2rem',
                            'fontWeight': '700'
                        })
                    ], style={'textAlign': 'center'})
                ], style={
                    'backgroundColor': 'white',
                    'padding': '30px 25px',
                    'borderRadius': '12px',
                    'boxShadow': '0 6px 20px rgba(220,53,69,0.15)',
                    'border': '1px solid #e9ecef',
                    'minWidth': '220px',
                    'transition': 'transform 0.3s ease, box-shadow 0.3s ease'
                }),
                
                # Conversion Rate Card
                html.Div([
                    html.Div([
                        html.I(className="fas fa-percentage", style={
                            'fontSize': '2rem', 
                            'color': '#ffc107',
                            'marginBottom': '15px'
                        }),
                        html.H4("Conversion Rate", style={
                            'color': '#333', 
                            'marginBottom': '12px',
                            'fontSize': '1.1rem',
                            'fontWeight': '600'
                        }),
                        html.H2(id="kpi-conversion-rate", style={
                            'color': '#ffc107', 
                            'margin': '0',
                            'fontSize': '2.2rem',
                            'fontWeight': '700'
                        })
                    ], style={'textAlign': 'center'})
                ], style={
                    'backgroundColor': 'white',
                    'padding': '30px 25px',
                    'borderRadius': '12px',
                    'boxShadow': '0 6px 20px rgba(255,193,7,0.15)',
                    'border': '1px solid #e9ecef',
                    'minWidth': '220px',
                    'transition': 'transform 0.3s ease, box-shadow 0.3s ease'
                })
                
            ], style={
                'display': 'flex', 
                'justifyContent': 'center', 
                'flexWrap': 'wrap', 
                'gap': '30px'
            }),
            
        ], style={'marginBottom': '50px'}),
        
        # Update the Analytics Chart Section in your app.layout
html.Div([
    html.H3("Analytics Overview", 
           style={
               'color': '#0051a6', 
               'textAlign': 'center', 
               'marginBottom': '30px',
               'fontSize': '1.6rem',
               'fontWeight': '600'
           }),
    html.Div([
        # First Chart - Total Sites Visited
        html.Div([
            dcc.Graph(
                id="age-chart",
                config={
                    'displayModeBar': True,
                    'displaylogo': False,
                    'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d']
                }
            )
        ], style={
            'backgroundColor': 'white',
            'borderRadius': '12px',
            'padding': '20px',
            'marginBottom': '30px',
            'boxShadow': '0 2px 8px rgba(0,0,0,0.1)'
        }),
        
        # Second Chart - Conversion Rate
        html.Div([
            dcc.Graph(
                id="conversion-chart",
                config={
                    'displayModeBar': True,
                    'displaylogo': False,
                    'modeBarButtonsToRemove': ['pan2d', 'lasso2d', 'select2d']
                }
            )
        ], style={
            'backgroundColor': 'white',
            'borderRadius': '12px',
            'padding': '20px',
            'boxShadow': '0 2px 8px rgba(0,0,0,0.1)'
        })
    ])
        ], style={
            'backgroundColor': 'white',
            'padding': '35px',
            'borderRadius': '15px',
            'boxShadow': '0 6px 20px rgba(0,0,0,0.1)',
            'border': '1px solid #e9ecef',
            'marginBottom': '40px'
        }),
        
        # Status Messages Section
        html.Div(id="status-message", style={
            'textAlign': 'center',
            'padding': '20px',
            'borderRadius': '8px',
            'marginTop': '20px'
        })
        
    ], style={
        'maxWidth': '1400px', 
        'margin': '0 auto', 
        'padding': '0 30px 40px 30px'
    })
    
], style={
    'backgroundColor': '#f5f7fa',
    'minHeight': '100vh',
    'fontFamily': '"Segoe UI", Tahoma, Geneva, Verdana, sans-serif'
})

# Updated callback with improved chart formatting
@app.callback(
    [Output('kpi-new-users', 'children'),
     Output('kpi-converted', 'children'),
     Output('kpi-conversion-rate', 'children'),
     Output('age-chart', 'figure'),
     Output('conversion-chart', 'figure')],
    [Input('age-category-dropdown', 'value')]
)
def update_dashboard(selected_age_categories):
    # Filter the dataframe based on selected age categories
    if 'all' in selected_age_categories or not selected_age_categories:
        df_filtered = df1
    else:
        # Handle multiple age category selections
        age_ranges = []
        for category in selected_age_categories:
            age_parts = category.split('-')
            min_age = int(age_parts[0])
            max_age = int(age_parts[1])
            age_ranges.append((min_age, max_age))
        
        # Create filter condition for all selected age ranges
        conditions = []
        for min_age, max_age in age_ranges:
            condition = (df1['age'] >= min_age) & (df1['age'] <= max_age)
            conditions.append(condition)
        
        # Combine conditions with OR logic
        if conditions:
            combined_condition = conditions[0]
            for condition in conditions[1:]:
                combined_condition = combined_condition | condition
            df_filtered = df1[combined_condition]
        else:
            df_filtered = df1
    
    # Calculate KPIs using the filtered data
    total_new_users, total_converted, conversion_rate = calculate_kpis(df_filtered)
    
    # Create both chart figures
    if not df_filtered.empty:
        # Create age bins for the charts
        min_age = (min(df_filtered['age']) // 5) * 5
        max_age = (max(df_filtered['age']) // 5) * 5
        age_bins = range(min_age, max_age + 6, 5)
        
        # Build age labels manually
        age_labels = []
        for i in age_bins[:-1]:
            label = str(i) + "-" + str(i + 4)
            age_labels.append(label)
        
        # Work with a copy of the filtered dataframe
        df_filtered_copy = df_filtered.copy()
        df_filtered_copy['AgeGroup'] = pd.cut(df_filtered_copy['age'], bins=age_bins, labels=age_labels, right=False)
        
        # First chart - Total Sites Visited
        age_group_stats = df_filtered_copy.groupby('AgeGroup')['total_pages_visited'].sum().reset_index()
        sites_figure = {
            'data': [{
                'x': age_group_stats['AgeGroup'].astype(str),
                'y': age_group_stats['total_pages_visited'],
                'type': 'bar',
                'marker': {
                    'color': '#0051a6',
                    'line': {'color': '#003d82', 'width': 1}
                },
                'hovertemplate': '<b>Age Group:</b> %{x}<br><b>Total Sites Visited:</b> %{y:,}<extra></extra>'
            }],
            'layout': {
                'title': {
                    'text': 'Total Sites Visited by Age Group',
                    'x': 0.5,
                    'font': {'size': 18, 'color': '#0051a6', 'family': 'Segoe UI'}
                },
                'xaxis': {
                    'title': {'text': 'Age Group', 'font': {'size': 14, 'color': '#333'}},
                    'tickfont': {'size': 12, 'color': '#666'},
                    'gridcolor': '#e9ecef'
                },
                'yaxis': {
                    'title': {'text': 'Total Sites Visited', 'font': {'size': 14, 'color': '#333'}},
                    'tickfont': {'size': 12, 'color': '#666'},
                    'gridcolor': '#e9ecef'
                },
                'plot_bgcolor': 'white',
                'paper_bgcolor': 'white',
                'font': {'family': 'Segoe UI'},
                'margin': {'l': 80, 'r': 40, 't': 80, 'b': 80},
                'hovermode': 'x'
            }
        }
        # Second chart - Conversion Rate
        conversion_stats = df_filtered_copy.groupby('AgeGroup')['converted'].mean().reset_index()
        conversion_figure = {
            'data': [{
                'x': conversion_stats['AgeGroup'].astype(str),
                'y': conversion_stats['converted'] * 100,  
                'type': 'bar',
                'marker': {
                    'color': '#28a745',
                    'line': {'color': '#218838', 'width': 1}
                },
                'hovertemplate': '<b>Age Group:</b> %{x}<br><b>Conversion Rate:</b> %{y:.1f}%<extra></extra>'
            }],
            'layout': {
                'title': {
                    'text': 'Conversion Rate by Age Group',
                    'x': 0.5,
                    'font': {'size': 18, 'color': '#0051a6', 'family': 'Segoe UI'}
                },
                'xaxis': {
                    'title': {'text': 'Age Group', 'font': {'size': 14, 'color': '#333'}},
                    'tickfont': {'size': 12, 'color': '#666'},
                    'gridcolor': '#e9ecef'
                },
                'yaxis': {
                    'title': {'text': 'Conversion Rate (%)', 'font': {'size': 14, 'color': '#333'}},
                    'tickfont': {'size': 12, 'color': '#666'},
                    'gridcolor': '#e9ecef',
                    'ticksuffix': '%'
                },
                'plot_bgcolor': 'white',
                'paper_bgcolor': 'white',
                'font': {'family': 'Segoe UI'},
                'margin': {'l': 80, 'r': 40, 't': 80, 'b': 80},
                'hovermode': 'x'
            }
        }
    else:
        # Empty figures if no data available
        empty_layout = {
            'title': {
                'text': 'No data available for selected filter',
                'x': 0.5,
                'font': {'size': 18, 'color': '#666', 'family': 'Segoe UI'}
            },
            'xaxis': {'title': {'text': 'Age Group', 'font': {'size': 14, 'color': '#333'}}},
            'yaxis': {'title': {'text': '', 'font': {'size': 14, 'color': '#333'}}},
            'plot_bgcolor': 'white',
            'paper_bgcolor': 'white',
            'font': {'family': 'Segoe UI'},
            'margin': {'l': 80, 'r': 40, 't': 80, 'b': 80}
        }
        sites_figure = {'data': [], 'layout': empty_layout.copy()}
        conversion_figure = {'data': [], 'layout': empty_layout.copy()}
        conversion_figure['layout']['yaxis']['title']['text'] = 'Conversion Rate (%)'
    
    # Return the formatted KPI values and chart figures
    total_users_formatted = str(total_new_users) + ","
    total_users_formatted = "{:,}".format(total_new_users)
    
    total_converted_formatted = "{:,}".format(total_converted)
    
    conversion_rate_formatted = str(conversion_rate) + "%"
    
    return (
        total_users_formatted, 
        total_converted_formatted, 
        conversion_rate_formatted, 
        sites_figure, 
        conversion_figure
    )

# Callback for PowerPoint download 
@app.callback(
    [Output("download-ppt", "data"),
     Output("status-message", "children")],
    [Input("download-btn", "n_clicks")],
    [State('age-category-dropdown', 'value')],
    prevent_initial_call=True
)
def download_ppt(n_clicks, selected_age_categories):
    if n_clicks > 0:
        try:
            # Filter the dataframe based on selected categories
            if 'all' in selected_age_categories or not selected_age_categories:
                df_filtered = df1
            else:
                # Handle multiple age category selections
                age_ranges = []
                for category in selected_age_categories:
                    age_parts = category.split('-')
                    min_age = int(age_parts[0])
                    max_age = int(age_parts[1])
                    age_ranges.append((min_age, max_age))
                
                # Create filter condition for all selected age ranges
                conditions = []
                for min_age, max_age in age_ranges:
                    condition = (df1['age'] >= min_age) & (df1['age'] <= max_age)
                    conditions.append(condition)
                
                # Combine conditions with OR logic
                if conditions:
                    combined_condition = conditions[0]
                    for condition in conditions[1:]:
                        combined_condition = combined_condition | condition
                    df_filtered = df1[combined_condition]
                else:
                    df_filtered = df1
            
            # Create presentation using template
            ppt_filename = create_presentation(df_filtered, template_path="Sales_presentation1.pptx")
            
            # Check if presentation was created successfully
            if ppt_filename and os.path.exists(ppt_filename):
                # Read the file and prepare for download
                with open(ppt_filename, 'rb') as f:
                    ppt_data = f.read()
                
                # Return success response
                success_message = html.Div("✅ Report downloaded successfully!", 
                       style={
                           'color': '#28a745', 
                           'fontWeight': 'bold',
                           'backgroundColor': '#d4edda',
                           'border': '1px solid #c3e6cb',
                           'padding': '12px 20px',
                           'borderRadius': '8px',
                           'display': 'inline-block'
                       })
                
                return (
                    dcc.send_bytes(ppt_data, filename=ppt_filename),
                    success_message
                )
            else:
                # Return error if presentation creation failed
                error_message = html.Div("Error generating report. Please check template file and debug output.", 
                       style={
                           'color': '#dc3545', 
                           'fontWeight': 'bold',
                           'backgroundColor': '#f8d7da',
                           'border': '1px solid #f5c6cb',
                           'padding': '12px 20px',
                           'borderRadius': '8px',
                           'display': 'inline-block'
                       })
                
                return (no_update, error_message)
                
        except Exception as e:
            # Handle any exceptions that occur
            import traceback
            error_details = traceback.format_exc()
            print("Download error: " + error_details)
            
            error_message = html.Div("Error: " + str(e), 
                   style={
                       'color': '#dc3545', 
                       'fontWeight': 'bold',
                       'backgroundColor': '#f8d7da',
                       'border': '1px solid #f5c6cb',
                       'padding': '12px 20px',
                       'borderRadius': '8px',
                       'display': 'inline-block'
                   })
            
            return (no_update, error_message)
    
    # Return nothing if button hasn't been clicked
    return no_update, ""


if __name__ == '__main__':
    app.run_server(debug=True, port=8070)